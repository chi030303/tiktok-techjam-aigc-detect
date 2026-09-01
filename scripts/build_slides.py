#!/usr/bin/env python3
# 2026-08-31, tianqi, contest 16:9 slides + figures for data/model/eval/badcase
"""Build 16:9 decks + PNG figures.

  MPLCONFIGDIR=/tmp python scripts/build_slides.py

Writes TechJam_Challenge5.pptx (14-page team) and TechJam_Challenge5_demo.pptx (6-page YouTube).
"""
from __future__ import annotations

import os
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", "/tmp/mplconfig")

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Rectangle
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap, qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
FIG = REPO / "docs" / "slides" / "figures"
OUT = REPO / "docs" / "slides" / "TechJam_Challenge5.pptx"
# 2026-08-31, tianqi, 6-slide deck for the 3–4 min YouTube demo
OUT_DEMO = REPO / "docs" / "slides" / "TechJam_Challenge5_demo.pptx"
# end

INK = "#0F172A"
MUTED = "#64748B"
ACCENT = "#E11D48"
TEAL = "#0F766E"
BLUE = "#1D4ED8"
AMBER = "#B45309"
LINE = "#E2E8F0"
CARD = "#F8FAFC"
OK = "#047857"
# end

W, H = 13.333, 7.5
RGB = {
    "ink": RGBColor(0x0F, 0x17, 0x2A),
    "muted": RGBColor(0x64, 0x74, 0x8B),
    "accent": RGBColor(0xE1, 0x1D, 0x48),
    "teal": RGBColor(0x0F, 0x76, 0x6E),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "card": RGBColor(0xF8, 0xFA, 0xFC),
    "line": RGBColor(0xE2, 0xE8, 0xF0),
}


def _setup_mpl() -> None:
    plt.rcParams.update(
        {
            "font.family": "sans-serif",
            "font.sans-serif": ["Arial", "Helvetica", "DejaVu Sans"],
            "font.size": 11,
            "axes.edgecolor": LINE,
            "axes.labelcolor": INK,
            "text.color": INK,
            "figure.facecolor": "white",
            "axes.facecolor": "white",
            "savefig.bbox": "tight",
            "savefig.dpi": 180,
        }
    )


def _round(ax, xy, w, h, fc, ec=LINE, lw=1.2, r=0.08, z=2):
    p = FancyBboxPatch(
        xy,
        w,
        h,
        boxstyle=f"round,pad=0.02,rounding_size={r}",
        facecolor=fc,
        edgecolor=ec,
        linewidth=lw,
        zorder=z,
    )
    ax.add_patch(p)
    return p


def _txt(ax, x, y, s, size=11, color=INK, ha="center", va="center", weight="normal"):
    ax.text(x, y, s, fontsize=size, color=color, ha=ha, va=va, weight=weight, zorder=4)


def fig_data_path() -> Path:
    fig, ax = plt.subplots(figsize=(12.4, 3.6))
    ax.set_xlim(0, 12.4)
    ax.set_ylim(0, 3.6)
    ax.axis("off")
    steps = [
        (0.2, "CIFAKE", "32×32 SD1.4\nairplane photos", "official 0.57", "#FEE2E2", ACCENT),
        (2.7, "SID 140k", "social FLUX-heavy\n+ online aug", "official 0.970", "#DBEAFE", BLUE),
        (5.2, "D3 mix", "UNet + ADM/DDPM\n+ flux2/sd35", "0.978 · Nova 0.988", "#CCFBF1", TEAL),
        (7.7, "D5 = D3+D4", "+ PixArt / SDXL\n+ GPT / nano", "0.975 · Nova 0.984", "#FEF3C7", AMBER),
        # 2026-09-01, tianqi, D6 finished: official 0.977, pair_acc 0.805, not a submit
        (10.2, "D6", "+ 118 i2i fakes\nCodex + nano", "0.977 · pair 0.81", "#EDE9FE", "#6D28D9"),
        # end
    ]
    for x, title, body, score, fc, ec in steps:
        _round(ax, (x, 0.55), 2.05, 2.55, fc, ec, 1.8, 0.12)
        _txt(ax, x + 1.02, 2.62, title, 13, INK, weight="bold")
        _txt(ax, x + 1.02, 1.72, body, 9, MUTED)
        _txt(ax, x + 1.02, 0.88, score, 10, ec, weight="bold")
    for x in (2.25, 5.15, 7.65, 10.15):
        ax.annotate(
            "",
            xy=(x + 0.42, 1.85),
            xytext=(x, 1.85),
            arrowprops=dict(arrowstyle="-|>", color=MUTED, lw=1.6),
        )
    _txt(ax, 6.2, 3.35, "Train mix path  (hold-out: official val DALL·E + EvalGEN never enter train)", 12, MUTED)
    dest = FIG / "data_path.png"
    fig.savefig(dest)
    plt.close()
    return dest


def fig_generators() -> Path:
    gens = [
        ("SID FLUX (kept)", 70000, BLUE),
        ("WildFake UNet", 3931, TEAL),
        ("flux2", 1500, TEAL),
        ("sd35", 1500, TEAL),
        ("nano t2i", 1500, AMBER),
        ("PixArt DiT", 1500, AMBER),
        ("SDXL", 1514, AMBER),
        ("GPT-image", 1500, AMBER),
        ("ADM pixel", 1000, "#0F766E"),
        ("DDPM pixel", 1000, "#0F766E"),
        ("i2i nano", 59, "#6D28D9"),
        ("i2i Codex", 59, "#6D28D9"),
    ]
    fig, ax = plt.subplots(figsize=(12.2, 4.4))
    names = [g[0] for g in gens]
    vals = [g[1] for g in gens]
    colors = [g[2] for g in gens]
    y = range(len(gens))[::-1]
    ax.barh(list(y), vals, color=colors, height=0.72, edgecolor="white")
    ax.set_yticks(list(y))
    ax.set_yticklabels(names, fontsize=10)
    ax.set_xlabel("images in train (SID FLUX is ~half of 140k; mixin counts are exact)")
    ax.set_xscale("log")
    ax.set_xlim(40, 120000)
    for yy, v in zip(y, vals):
        ax.text(v * 1.08, yy, f"{v:,}", va="center", fontsize=9, color=INK)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.set_title("Generator coverage after D6  ·  architecture mix, not “newer T2I”", loc="left", fontsize=13, pad=8)
    dest = FIG / "generators.png"
    fig.savefig(dest)
    plt.close()
    return dest


def fig_decisions() -> Path:
    fig, ax = plt.subplots(figsize=(12.4, 3.8))
    ax.set_xlim(0, 12.4)
    ax.set_ylim(0, 3.8)
    ax.axis("off")
    cards = [
        (0.25, "Why not GAN as train", "Contest val = DALL·E.\nEvalGEN = flow / AR / DiT.\nB-axis skipped: leave GAN\nas unseen, not a 10/20% mix.\nPixel hole filled by ADM/DDPM\n(not StyleGAN)."),
        (4.3, "i2i is triplets, not 8k", "59 complete scenes:\nreal + nano + Codex.\nA2 i2i-only official 0.443,\npair_acc 0.42 (overfit).\nD6 = D5 + 118 fakes:\nofficial 0.977, pair 0.805."),
        (8.35, "No SID tampered as i2i", "SID label=2 is local edit.\nContest is image-level AIGC.\nDo not treat tampered as\nwhole-image i2i. Do not mine\nDALL·E / EvalGEN into train."),
    ]
    for x, title, body in cards:
        _round(ax, (x, 0.25), 3.8, 3.3, CARD, LINE, 1.2, 0.1)
        _txt(ax, x + 1.9, 3.15, title, 13, INK, weight="bold")
        _txt(ax, x + 1.9, 1.55, body, 10, MUTED)
    dest = FIG / "decisions.png"
    fig.savefig(dest)
    plt.close()
    return dest


# 2026-09-01, tianqi, A-grid + D6 i2i numbers for team slides
def fig_i2i_table() -> Path:
    fig, ax = plt.subplots(figsize=(12.4, 3.6))
    ax.axis("off")
    cols = ["Ablation", "Train", "Official 400", "EvalGEN clean", "i2i pair_acc", "Takeaway"]
    rows = [
        ["A1 t2i-only", "t2i, no SID", "0.595", "0.760", "0.653", "no social domain"],
        ["A2 59-triplet i2i", "118 fakes only", "0.443", "0.810", "0.424", "overfits 59 scenes"],
        ["A3 t2i + 118 i2i", "still no SID", "0.786", "0.786", "0.593", "i2i cannot replace SID"],
        ["D6 = D5 + 118 i2i", "SID mix + 118 fakes", "0.977", "0.994", "0.805", "pair up; contest unchanged"],
    ]
    table = ax.table(
        cellText=rows,
        colLabels=cols,
        loc="center",
        cellLoc="center",
        colColours=["#F1F5F9"] * 6,
    )
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    table.scale(1.18, 1.55)
    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#E2E8F0")
        if r == 0:
            cell.set_text_props(weight="bold", color=INK)
        if r == 2:
            cell.set_facecolor("#FEE2E2")
        if r == 4:
            cell.set_facecolor("#EDE9FE")
    ax.set_title(
        "i2i pair_acc = P(recon_score > paired real) on 118 pairs  ·  A-grid has no SID  ·  D6 keeps SID",
        loc="left",
        fontsize=12,
        color=MUTED,
        pad=8,
    )
    dest = FIG / "i2i_table.png"
    fig.savefig(dest)
    plt.close()
    return dest
# end


def fig_transforms() -> Path:
    fig, ax = plt.subplots(figsize=(12.4, 3.5))
    ax.set_xlim(0, 12.4)
    ax.set_ylim(0, 3.5)
    ax.axis("off")
    ops = [
        ("JPEG", "q90 / 70 / 50 / 30", "weakest: q30"),
        ("Blur", "σ 0.5 / 1.0 / 2.0", "social re-encode"),
        ("Resize", "×0.5 / ×0.25 then up", "weakest: ×0.25"),
        ("Noise", "σ 0.02 / 0.05 / 0.10", "[0,1] pixels"),
        ("Jitter", "±20% B/C/S", "one sample / image"),
        ("Crop", "center 80%, no upscale", "small thumbnail"),
    ]
    for i, (t, body, note) in enumerate(ops):
        x = 0.25 + i * 2.02
        _round(ax, (x, 0.7), 1.88, 2.15, CARD, LINE, 1.1, 0.1)
        _txt(ax, x + 0.94, 2.45, t, 13, TEAL, weight="bold")
        _txt(ax, x + 0.94, 1.75, body, 9, MUTED)
        _txt(ax, x + 0.94, 1.05, note, 8, AMBER)
    _txt(
        ax,
        6.2,
        3.2,
        "14 frozen eval keys + clean  ·  same ops as online train aug (p_clean=0.2)  ·  robust AUC = mean of 14",
        11,
        MUTED,
    )
    dest = FIG / "transforms.png"
    fig.savefig(dest)
    plt.close()
    return dest


def fig_model_bars() -> Path:
    labels = [
        "CIFAKE CLIP-B",
        "C-Pixel FFT/ADM",
        "SID frozen",
        "SID consistency",
        "SID dual RGB+HP",
        "SID 336",
        "first-4",
        "CLIP-L last-4",
        "D3 frozen",
        "D3 dualbranch",
        "CLIP-B last-4",
        "fuse last4+D3",
    ]
    vals = [0.569, 0.649, 0.970, 0.965, 0.966, 0.969, 0.974, 0.980, 0.978, 0.983, 0.990, 0.993]
    colors = [ACCENT] * 2 + [MUTED] * 6 + [TEAL, TEAL, BLUE, ACCENT]
    fig, ax = plt.subplots(figsize=(12.4, 4.5))
    bars = ax.barh(labels[::-1], vals[::-1], color=colors[::-1], height=0.68, edgecolor="white")
    ax.set_xlim(0.50, 1.0)
    ax.set_xlabel("official val 400 formula  (0.50×AUC_clean + 0.50×AUC_robust)")
    ax.axvline(0.990, color=BLUE, ls="--", lw=1, alpha=0.7)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for b, v in zip(bars, vals[::-1]):
        ax.text(v + 0.004, b.get_y() + b.get_height() / 2, f"{v:.3f}", va="center", fontsize=9)
    ax.set_title("Model path: MLP/FFT and extra towers lose to last-4 ranking + D3 fuse", loc="left", fontsize=13)
    dest = FIG / "model_bars.png"
    fig.savefig(dest)
    plt.close()
    return dest


def fig_last4_fuse() -> Path:
    fig, ax = plt.subplots(figsize=(12.4, 3.9))
    ax.set_xlim(0, 12.4)
    ax.set_ylim(0, 3.9)
    ax.axis("off")
    _txt(ax, 3.1, 3.55, "CLIP-B/16 last-4  (submit if 1 ckpt)", 13, BLUE, weight="bold")
    for i in range(12):
        x = 0.25 + i * 0.48
        fc = "#93C5FD" if i >= 8 else "#E2E8F0"
        _round(ax, (x, 2.15), 0.42, 0.95, fc, LINE, 0.8, 0.06)
        _txt(ax, x + 0.21, 2.62, str(i + 1), 8, INK)
    _txt(ax, 2.15, 1.95, "frozen 1–8", 9, MUTED)
    _txt(ax, 5.15, 1.95, "unfreeze 9–12", 9, BLUE, weight="bold")
    _round(ax, (6.15, 2.25), 1.55, 0.75, "#DBEAFE", BLUE, 1.4, 0.08)
    _txt(ax, 6.92, 2.62, "linear", 11, BLUE, weight="bold")
    _txt(ax, 9.7, 3.55, "Mean-logit fuse  (if 2 ckpts)", 13, ACCENT, weight="bold")
    _round(ax, (8.35, 2.35), 1.7, 0.85, "#DBEAFE", BLUE, 1.2, 0.08)
    _txt(ax, 9.2, 2.77, "last4 logit", 10, BLUE)
    _round(ax, (8.35, 1.15), 1.7, 0.85, "#CCFBF1", TEAL, 1.2, 0.08)
    _txt(ax, 9.2, 1.57, "D3 mix logit", 10, TEAL)
    _round(ax, (10.3, 1.7), 1.85, 0.95, "#FEE2E2", ACCENT, 1.6, 0.1)
    _txt(ax, 11.22, 2.17, "mean → σ\n0.993", 11, ACCENT, weight="bold")
    ax.annotate("", xy=(10.3, 2.55), xytext=(10.05, 2.75), arrowprops=dict(arrowstyle="-|>", color=MUTED, lw=1.3))
    ax.annotate("", xy=(10.3, 1.8), xytext=(10.05, 1.55), arrowprops=dict(arrowstyle="-|>", color=MUTED, lw=1.3))
    # 2026-09-01, tianqi, D5/D6 fuse did not beat D3 fuse
    _txt(
        ax,
        6.2,
        0.45,
        "Do not train last-4 on the mix → 0.976.  D5/D6 fuse 0.9927 / 0.9929 do not beat D3 fuse 0.993.",
        11,
        MUTED,
    )
    # end
    dest = FIG / "last4_fuse.png"
    fig.savefig(dest)
    plt.close()
    return dest


# 2026-09-01, tianqi, why fuse beats stacking / extra mix-in
def fig_fuse_why() -> Path:
    fig, ax = plt.subplots(figsize=(12.4, 2.6))
    ax.axis("off")
    cols = ["", "last-4", "D3 mix", "fuse last4+D3", "train last-4 on mix"]
    rows = [
        ["Official 400 formula", "0.990", "0.978", "0.993", "0.976"],
        ["Nova recall @0.5", "0.49", "0.86", "0.56", "—"],
        ["Nova 15-cond formula", "—", "0.9865", "0.9880", "—"],
        ["400 @0.5  FP / FN", "1 / 60", "4 / 20", "1 / 44", "—"],
    ]
    table = ax.table(
        cellText=rows,
        colLabels=cols,
        loc="center",
        cellLoc="center",
        colColours=["#F1F5F9"] * 5,
    )
    table.auto_set_font_size(False)
    table.set_fontsize(11)
    table.scale(1.18, 1.45)
    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#E2E8F0")
        if r == 0:
            cell.set_text_props(weight="bold", color=INK)
        if c == 3 and r > 0:
            cell.set_facecolor("#FEE2E2")
            cell.set_text_props(weight="bold")
        if c == 4 and r == 1:
            cell.set_facecolor("#FEF3C7")
    ax.set_title(
        "Fuse at inference  ·  last-4 Nova 15-cond not rerun  ·  D5/D6 fuse 0.9927 / 0.9929 vs D3 0.993",
        loc="left",
        fontsize=11,
        color=MUTED,
        pad=6,
    )
    dest = FIG / "fuse_why.png"
    fig.savefig(dest)
    plt.close()
    return dest
# end


def fig_fuse_arch() -> Path:
    # 2026-08-31, tianqi, demo fuse diagram: no 1-ckpt / 2-ckpt labels, no stack caption
    fig, ax = plt.subplots(figsize=(12.4, 3.5))
    ax.set_xlim(0, 12.4)
    ax.set_ylim(0, 3.5)
    ax.axis("off")
    _txt(ax, 3.1, 3.2, "CLIP-B/16 last-4", 14, BLUE, weight="bold")
    for i in range(12):
        x = 0.25 + i * 0.48
        fc = "#93C5FD" if i >= 8 else "#E2E8F0"
        _round(ax, (x, 1.85), 0.42, 0.95, fc, LINE, 0.8, 0.06)
        _txt(ax, x + 0.21, 2.32, str(i + 1), 8, INK)
    _txt(ax, 2.15, 1.65, "frozen 1–8", 9, MUTED)
    _txt(ax, 5.15, 1.65, "unfreeze 9–12", 9, BLUE, weight="bold")
    _round(ax, (6.15, 1.95), 1.55, 0.75, "#DBEAFE", BLUE, 1.4, 0.08)
    _txt(ax, 6.92, 2.32, "linear", 11, BLUE, weight="bold")
    _txt(ax, 9.7, 3.2, "Mean-logit fuse", 14, ACCENT, weight="bold")
    _round(ax, (8.35, 2.05), 1.7, 0.8, "#DBEAFE", BLUE, 1.2, 0.08)
    _txt(ax, 9.2, 2.45, "last-4 logit", 10, BLUE)
    _round(ax, (8.35, 0.95), 1.7, 0.8, "#CCFBF1", TEAL, 1.2, 0.08)
    _txt(ax, 9.2, 1.35, "mixed-data logit", 10, TEAL)
    _round(ax, (10.3, 1.45), 1.85, 0.95, "#FEE2E2", ACCENT, 1.6, 0.1)
    _txt(ax, 11.22, 1.92, "mean → σ\n0.993", 11, ACCENT, weight="bold")
    ax.annotate("", xy=(10.3, 2.25), xytext=(10.05, 2.45), arrowprops=dict(arrowstyle="-|>", color=MUTED, lw=1.3))
    ax.annotate("", xy=(10.3, 1.6), xytext=(10.05, 1.35), arrowprops=dict(arrowstyle="-|>", color=MUTED, lw=1.3))
    dest = FIG / "fuse_arch.png"
    fig.savefig(dest)
    plt.close()
    return dest
    # end


def fig_arch_table() -> Path:
    # 2026-09-01, tianqi, demo mix: SID last-4 vs mixed-head families vs generated-not-submit
    fig, ax = plt.subplots(figsize=(12.4, 4.05))
    ax.set_xlim(0, 12.4)
    ax.set_ylim(0, 4.05)
    ax.axis("off")
    _round(ax, (0.15, 3.28), 12.1, 0.68, "#CCFBF1", TEAL, 1.6, 0.1)
    _txt(ax, 6.2, 3.74, "last-4 (1 ckpt)  ·  SID_Set  ~140,000", 14, TEAL, weight="bold", ha="center")
    _txt(
        ax,
        6.2,
        3.42,
        "≈70k real + ~70k FLUX.  Mixed head replaces a thin slice of SID FLUX — not stacked, not every family below.",
        10,
        INK,
        ha="center",
    )
    _txt(ax, 0.22, 3.08, "In the mixed head  (fuse 2nd ckpt)", 11, TEAL, ha="left", weight="bold")
    in_head = [
        ("WildFake original SD", "3,931"),
        ("Flux.2", "1,500"),
        ("SD3.5 Medium", "1,500"),
        ("ADM", "1,000"),
        ("DDPM", "1,000"),
    ]
    for i, (name, n) in enumerate(in_head):
        x = 0.15 + i * 2.42
        _round(ax, (x, 2.18), 2.30, 0.78, "#CCFBF1", TEAL, 1.0, 0.08)
        _txt(ax, x + 0.14, 2.72, name, 9, MUTED, ha="left")
        _txt(ax, x + 0.14, 2.38, n, 14, INK, ha="left", weight="bold")
    _txt(
        ax,
        0.22,
        1.95,
        "Also generated  ·  extra T2I did not lift contest score  ·  not in submit",
        11,
        AMBER,
        ha="left",
        weight="bold",
    )
    built = [
        ("PixArt-Sigma", "1,500"),
        ("SDXL", "1,514"),
        ("GPT-image", "1,500"),
        ("nano banana (t2i)", "1,500"),
    ]
    for i, (name, n) in enumerate(built):
        x = 0.15 + i * 3.05
        _round(ax, (x, 0.95), 2.90, 0.78, "#F8FAFC", LINE, 0.9, 0.08)
        _txt(ax, x + 0.16, 1.50, name, 10, MUTED, ha="left")
        _txt(ax, x + 0.16, 1.16, n, 14, MUTED, ha="left", weight="bold")
    _txt(ax, 6.2, 0.55, "Self-built from local checkpoints + ComfyUI  ·  will open-source  ·  no GAN / Hunyuan", 10, MUTED)
    dest = FIG / "arch_table.png"
    fig.savefig(dest, bbox_inches="tight", pad_inches=0.08)
    plt.close()
    return dest
    # end


def fig_eval_full() -> Path:
    import numpy as np

    keys = ["clean", "J90", "J70", "J50", "J30", "b0.5", "b1", "b2", "×0.5", "×0.25", "n02", "n05", "n10", "jit", "crop"]
    sid_full = [0.9655, 0.9683, 0.9630, 0.9622, 0.9566, 0.9701, 0.9763, 0.9722, 0.9752, 0.9641, 0.9572, 0.9666, 0.9726, 0.9626, 0.9548]
    last4 = [0.991, 0.991, 0.988, 0.987, 0.983, 0.992, 0.993, 0.985, 0.991, 0.980, 0.991, 0.987, 0.988, 0.990, 0.991]
    mixed = [0.985, 0.984, 0.978, 0.954, 0.944, 0.985, 0.986, 0.973, 0.984, 0.958, 0.974, 0.969, 0.967, 0.980, 0.971]
    fuse = [0.995, 0.995, 0.992, 0.988, 0.984, 0.996, 0.996, 0.989, 0.994, 0.984, 0.992, 0.989, 0.989, 0.994, 0.993]
    fig, (ax, axb) = plt.subplots(1, 2, figsize=(12.4, 4.15), gridspec_kw={"width_ratios": [2.15, 1]})
    x = range(15)
    ax.plot(x, fuse, "o-", color=ACCENT, lw=2.1, label="fuse  0.993  (n=400)", ms=4)
    ax.plot(x, last4, "s-", color=BLUE, lw=1.7, label="last-4  0.990 / full 0.989", ms=4)
    ax.plot(x, mixed, "^-", color=TEAL, lw=1.6, label="mixed data  0.978  (n=400)", ms=4)
    ax.plot(x, sid_full, "x--", color=MUTED, lw=1.3, label="SID  0.966  (full 13,843)", ms=4)
    ax.set_xticks(list(x))
    ax.set_xticklabels(keys, fontsize=8)
    ax.set_ylim(0.93, 1.0)
    ax.set_ylabel("AUROC")
    ax.set_title("Official val  (COCO + DALL·E, never trained)", loc="left", fontsize=11)
    ax.legend(frameon=False, fontsize=8, loc="lower left")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    labs = ["EvalGEN", "Nova AUC", "Nova rec@0.5"]
    last4b = [0.989, 0.963, 0.49]
    mixedb = [0.995, 0.988, 0.86]
    fuseb = [0.997, 0.988, 0.56]
    xx = np.arange(len(labs))
    w = 0.25
    axb.bar(xx - w, last4b, w, color=BLUE, label="last-4", edgecolor="white")
    axb.bar(xx, mixedb, w, color=TEAL, label="mixed data", edgecolor="white")
    axb.bar(xx + w, fuseb, w, color=ACCENT, label="fuse", edgecolor="white")
    axb.set_xticks(xx)
    axb.set_xticklabels(labs, fontsize=9)
    axb.set_ylim(0, 1.08)
    axb.set_title("EvalGEN hold-out  (unseen)", loc="left", fontsize=11)
    axb.legend(frameon=False, fontsize=8)
    axb.spines["top"].set_visible(False)
    axb.spines["right"].set_visible(False)
    fig.tight_layout()
    dest = FIG / "eval_full.png"
    fig.savefig(dest)
    plt.close()
    return dest


def fig_15cond_table() -> Path:
    cols = ["clean", "J90", "J70", "J50", "J30", "b0.5", "b1", "b2", "x0.5", "x0.25", "n02", "n05", "n10", "jit", "crop"]
    models = [
        ("fuse", [0.995, 0.995, 0.992, 0.988, 0.984, 0.996, 0.996, 0.989, 0.994, 0.984, 0.992, 0.989, 0.989, 0.994, 0.993]),
        ("last-4", [0.991, 0.991, 0.988, 0.987, 0.983, 0.992, 0.993, 0.985, 0.991, 0.980, 0.991, 0.987, 0.988, 0.990, 0.991]),
        ("mixed data", [0.985, 0.984, 0.978, 0.954, 0.944, 0.985, 0.986, 0.973, 0.984, 0.958, 0.974, 0.969, 0.967, 0.980, 0.971]),
        ("SID", [0.969, 0.971, 0.968, 0.968, 0.964, 0.972, 0.978, 0.978, 0.979, 0.973, 0.961, 0.971, 0.976, 0.963, 0.955]),
    ]
    cell = [[m] + [f"{v:.3f}" for v in row] for m, row in models]
    fig, ax = plt.subplots(figsize=(12.6, 2.6))
    ax.axis("off")
    table = ax.table(cellText=cell, colLabels=["model"] + cols, loc="center", cellLoc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(8)
    table.scale(1.05, 1.85)
    for (r, c), cell_ in table.get_celld().items():
        cell_.set_edgecolor(LINE)
        cell_.set_linewidth(0.5)
        if r == 0:
            cell_.set_facecolor("#0F172A")
            cell_.set_text_props(color="white", weight="bold")
        elif r == 1:
            cell_.set_facecolor("#FEE2E2")
        else:
            cell_.set_facecolor("white" if r % 2 else "#F8FAFC")
    ax.set_title("Official val AUROC  ·  15 conditions  ·  n=400 screen  ·  score = 0.50 clean + 0.50 mean of 14", loc="left", fontsize=11, pad=8)
    dest = FIG / "table_15cond.png"
    fig.savefig(dest, bbox_inches="tight", pad_inches=0.12)
    plt.close()
    return dest


def fig_eval_official() -> Path:
    keys = [
        "clean",
        "J90",
        "J70",
        "J50",
        "J30",
        "b0.5",
        "b1",
        "b2",
        "×0.5",
        "×0.25",
        "n02",
        "n05",
        "n10",
        "jit",
        "crop",
    ]
    fuse = [0.995, 0.995, 0.992, 0.988, 0.984, 0.996, 0.996, 0.989, 0.994, 0.984, 0.992, 0.989, 0.989, 0.994, 0.993]
    last4 = [0.991, 0.991, 0.988, 0.987, 0.983, 0.992, 0.993, 0.985, 0.991, 0.980, 0.991, 0.987, 0.988, 0.990, 0.991]
    d3 = [0.985, 0.984, 0.978, 0.954, 0.944, 0.985, 0.986, 0.973, 0.984, 0.958, 0.974, 0.969, 0.967, 0.980, 0.971]
    cifake = [0.561, 0.572, 0.582, 0.619, 0.654, 0.571, 0.567, 0.520, 0.552, 0.528, 0.583, 0.588, 0.613, 0.557, 0.587]
    x = range(15)
    fig, ax = plt.subplots(figsize=(12.4, 4.3))
    ax.plot(x, fuse, "o-", color=ACCENT, lw=2.2, label="fuse last4+D3  0.993", ms=5)
    ax.plot(x, last4, "s-", color=BLUE, lw=1.8, label="CLIP-B last-4  0.990", ms=5)
    ax.plot(x, d3, "^-", color=TEAL, lw=1.6, label="D3 mix  0.978", ms=5)
    ax.plot(x, cifake, "x--", color=MUTED, lw=1.2, label="CIFAKE CLIP-B  0.57", ms=5)
    ax.set_xticks(list(x))
    ax.set_xticklabels(keys, fontsize=9)
    ax.set_ylim(0.50, 1.01)
    ax.set_ylabel("AUROC on official val 400")
    ax.legend(frameon=False, loc="lower right")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.set_title("Official demonstration val  ·  200 COCO real + 200 DALL·E Advanced  ·  never trained", loc="left")
    dest = FIG / "eval_official.png"
    fig.savefig(dest)
    plt.close()
    return dest


def fig_evalgen() -> Path:
    import numpy as np

    labs = ["DALL·E\n(official)", "EvalGEN\nall", "Nova", "Infinity", "Nova rec@0.5"]
    last4 = [0.991, 0.989, 0.963, 0.983, 0.49]
    d3 = [0.985, 0.995, 0.988, 0.994, 0.86]
    fuse = [0.995, 0.997, 0.988, 0.995, 0.56]
    x = np.arange(len(labs))
    w = 0.25
    fig, ax = plt.subplots(figsize=(12.2, 4.3))
    ax.bar(x - w, last4, w, color=BLUE, label="last-4", edgecolor="white")
    ax.bar(x, d3, w, color=TEAL, label="D3 mix", edgecolor="white")
    ax.bar(x + w, fuse, w, color=ACCENT, label="fuse", edgecolor="white")
    ax.set_xticks(x)
    ax.set_xticklabels(labs)
    ax.set_ylim(0, 1.08)
    ax.set_ylabel("AUROC  (last group = recall at 0.5)")
    ax.legend(frameon=False, ncol=3, loc="upper left")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.set_title("Hold-out generators  ·  EvalGEN never in train  ·  Nova is the hard family", loc="left")
    dest = FIG / "evalgen.png"
    fig.savefig(dest)
    plt.close()
    return dest


def _font(size: int, bold: bool = False):
    for name in (
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        "/Library/Fonts/Arial.ttf",
    ):
        p = Path(name)
        if p.is_file():
            try:
                return ImageFont.truetype(str(p), size)
            except OSError:
                continue
    return ImageFont.load_default()


def fig_gallery() -> Path:
    img = Image.new("RGB", (1480, 860), (248, 250, 252))
    d = ImageDraw.Draw(img)
    font_lg = _font(28, True)
    font = _font(18)
    font_sm = _font(15)
    d.rectangle((0, 0, 1480, 72), fill=(15, 23, 42))
    d.text((28, 20), "Bad-case gallery   fuse last4+D3   official val 400  clean   threshold 0.5", fill="white", font=font_lg)
    d.text((28, 92), "How to read:  FP = real accused as AIGC (expensive).  FN = AIGC scored as real.  Sorted by worst confidence.  Cards carry generator / condition / path.", fill=(100, 116, 139), font=font)

    d.rectangle((28, 130, 720, 820), fill="white", outline=(226, 232, 240), width=2)
    d.rectangle((28, 130, 720, 178), fill=(254, 226, 226))
    d.text((44, 142), "FP  ·  1 / 200 COCO   (highest pred first)", fill=(185, 28, 28), font=font)

    d.rounded_rectangle((48, 200, 700, 430), 12, fill=(255, 255, 255), outline=(252, 165, 165), width=3)
    d.rectangle((60, 214, 250, 404), fill=(226, 232, 240))
    d.text((78, 290), "COCO photo", fill=(71, 85, 105), font=font_sm)
    d.text((268, 220), "pred 0.62   label REAL", fill=(185, 28, 28), font=font)
    d.text((268, 258), "Street / indoor photograph. Almost no FPs:\nlast-4 and fuse are low false-accusation\nsystems. This one FP is the costly error\nfor a platform (accusing a real user).", fill=(51, 65, 85), font=font_sm)

    d.rectangle((760, 130, 1452, 820), fill="white", outline=(226, 232, 240), width=2)
    d.rectangle((760, 130, 1452, 178), fill=(255, 237, 213))
    d.text((776, 142), "FN  ·  44 / 200 DALL·E   (lowest pred first)   many non-photoreal styles", fill=(180, 83, 9), font=font)

    styles = [
        ((780, 200), "#FDE68A", "illustration"),
        ((1118, 200), "#FBCFE8", "3D / cartoon"),
        ((780, 500), "#DDD6FE", "painterly"),
        ((1118, 500), "#E2E8F0", "photoreal miss"),
    ]
    notes = [
        "Poster / graphic DALL·E.\nLooks nothing like COCO.\nStill below 0.5 on last-4.",
        "Stylized character render.\nSID is social-photo FLUX;\nstyle shift → under-score.",
        "Oil / watercolor look.\nNon-photoreal cluster of FN\nin the 400 gallery.",
        "Hard ranking case:\nlooks like a real photo.\npred very low — model is sure.",
    ]
    for (xy, col, title), note in zip(styles, notes):
        x, y = xy
        d.rounded_rectangle((x, y, x + 310, y + 270), 12, outline=(253, 186, 116), width=2, fill="white")
        d.rectangle((x + 12, y + 12, x + 298, y + 118), fill=tuple(int(col[i : i + 2], 16) for i in (1, 3, 5)))
        d.text((x + 24, y + 50), title, fill=(15, 23, 42), font=font)
        d.text((x + 16, y + 136), note, fill=(71, 85, 105), font=font_sm)

    dest = FIG / "gallery.png"
    img.save(dest, "PNG")
    return dest


def fig_error_loop() -> Path:
    fig, ax = plt.subplots(figsize=(12.4, 3.2))
    ax.set_xlim(0, 12.4)
    ax.set_ylim(0, 3.2)
    ax.axis("off")
    rows = [
        (0.2, ACCENT, "CIFAKE ~random on official", "→ SID 140k + online 14 keys", "domain, not a bigger MLP"),
        (4.3, TEAL, "FLUX head, Nova rec 0.49", "→ D3 UNet/ADM/DDPM mix", "D4 new T2I did not lift Nova"),
        (8.4, BLUE, "last4 FN 60 vs D3 FN 20", "→ mean-logit fuse, no retrain", "FN 60→44, FP stays 1"),
    ]
    for x, c, a, b, n in rows:
        _round(ax, (x, 0.35), 3.85, 2.5, CARD, c, 1.6, 0.1)
        _txt(ax, x + 1.92, 2.4, a, 11, INK, weight="bold")
        _txt(ax, x + 1.92, 1.55, b, 12, c, weight="bold")
        _txt(ax, x + 1.92, 0.75, n, 10, MUTED)
    dest = FIG / "error_loop.png"
    fig.savefig(dest)
    plt.close()
    return dest


def _set_run(run, size=18, color=None, bold=False, font="Arial"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    if color is not None:
        run.font.color.rgb = color


def _add_title(slide, text, y=0.28, size=28, color="ink"):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(y), Inches(12.4), Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run(run, size, RGB[color], bold=True)
    return box


def _add_sub(slide, text, y=0.78, size=14):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(y), Inches(12.4), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run(run, size, RGB["muted"], False)
    return box


def _enable_bullet(paragraph) -> None:
    # 2026-08-31, tianqi, real PPT bullets on body copy (not a dash prefix)
    pPr = paragraph._p.get_or_add_pPr()
    for child in list(pPr):
        if child.tag.endswith("}buNone") or child.tag.endswith("}buChar") or child.tag.endswith("}buFont"):
            pPr.remove(child)
    buFont = OxmlElement("a:buFont")
    buFont.set("typeface", "Arial")
    buChar = OxmlElement("a:buChar")
    buChar.set("char", "•")
    pPr.append(buFont)
    pPr.append(buChar)
    pPr.set("marL", str(int(Inches(0.28))))
    pPr.set("indent", str(int(Inches(-0.18))))
    paragraph.level = 0
    # end


def _bullets(slide, lines, left=0.5, top=1.25, width=12.3, height=5.8, size=18, bullet=True):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = line.lstrip("- ").lstrip("• ").strip()
        _set_run(run, size, RGB["ink"], False)
        if bullet:
            _enable_bullet(p)
    return box


def _caption(slide, text, left, top, width=6.2, size=13):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.32))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run(run, size, RGB["ink"], True)
    return box


def _notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def _pic(slide, path: Path, top=1.15, height=5.9):
    slide.shapes.add_picture(str(path), Inches(0.4), Inches(top), width=Inches(12.5), height=Inches(height))


def _pic_at(slide, path: Path, left: float, top: float, width: float):
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def _accent_bar(slide) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGB["accent"]
    bar.line.fill.background()


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _footer(slide, n, total=14):
    box = slide.shapes.add_textbox(Inches(11.6), Inches(7.15), Inches(1.4), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{n} / {total}"
    _set_run(run, 11, RGB["muted"])


def build_pptx(figs: dict[str, Path]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 1 title
    s = prs.slides.add_slide(blank)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGB["accent"]
    bar.line.fill.background()
    _add_title(s, "Robust AIGC detection under JPEG / blur / crop", 2.2, 32)
    _add_sub(s, "TikTok TechJam 2026  ·  Challenge 5  ·  CLIP-B last-4  0.990   ·  fuse last4+D3  0.993", 3.0, 18)
    _bullets(
        s,
        ["Score = 0.50×AUC_clean + 0.50×AUC_robust   (not Acc@0.5)", "Model ≪ 2B   ·  official val and EvalGEN never trained"],
        0.5,
        4.0,
        12,
        1.6,
        16,
    )
    _notes(s, "开场：社交流转后的图会被压、糊、裁。分数是干净 AUC 和 14 档变换 AUC 各一半，不是 0.5 准确率。提交 last-4 或双权重 fuse。")
    _footer(s, 1)

    # 2 agenda
    s = prs.slides.add_slide(blank)
    _add_title(s, "Four sections")
    _add_sub(s, "Each section is one story, not a dump of ablations.")
    _bullets(
        s,
        [
            "1  Data — CIFAKE → SID → D3 → D5 → D6; generators, i2i, no-GAN, transforms",
            "2  Model — MLP/FFT → unfreeze / consistency / RGB+freq / 336 → why fuse",
            "3  Eval — official demonstration val, then held-out EvalGEN (Nova)",
            "4  Bad cases — FN cluster in non-photoreal styles + the HTML gallery",
        ],
        size=20,
    )
    _notes(s, "四页主线。数据讲分布和抉择；模型讲为什么最后 fuse；评估分官方和 EvalGEN；badcase 用画册。")
    _footer(s, 2)

    # 3 data path
    s = prs.slides.add_slide(blank)
    _add_title(s, "Data path: domain first, then architecture coverage")
    _add_sub(s, "CIFAKE 32×32 does not transfer. SID is social. Mix-ins replace equal SID FLUX.")
    _pic(s, figs["data_path"], 1.15, 3.55)
    # 2026-09-01, tianqi, D6 finished: official 0.977 pair 0.805, fuse 0.9929 no submit change
    _bullets(
        s,
        [
            "D3 is the mix that lifts Nova. D5 = D3 ∪ D4 ≈ D3 on official (0.975 vs 0.978), not a jump.",
            "D6 = D5 + 118 whole-image i2i fakes (no paired reals): official 0.977, pair_acc 0.805. Fuse last4+D6 = 0.9929 vs D3 fuse 0.9930 — do not submit D6.",
        ],
        0.5,
        4.85,
        12.3,
        1.8,
        16,
    )
    _notes(
        s,
        "CIFAKE 官方 0.57，换 SID 社交域加 online aug 到 0.970。D3 补 UNet 和 pixel，Nova 起来。D5 并上 PixArt/GPT 官方还略掉。D6 只加 118 张 i2i fake：官方 0.977，pair_acc 0.805，fuse 0.9929，仍提交 last-4 或 last4+D3。",
    )
    # end
    _footer(s, 3)

    # 4 generators
    s = prs.slides.add_slide(blank)
    _add_title(s, "Generator distribution: fill holes, do not stack more FLUX")
    _pic(s, figs["generators"], 1.05, 5.7)
    _notes(
        s,
        "SID 已经是 FLUX 海。mixin 是等权换掉同等数量 FLUX。D3 的 WildFake UNet 4k 和 ADM/DDPM 1k 才是 Nova 的活性成分。D4 的新 T2I 数量看起来均衡，但对 Nova 几乎没贡献。",
    )
    _footer(s, 4)

    # 5 decisions + i2i grid
    s = prs.slides.add_slide(blank)
    _add_title(s, "i2i: 59 triplets, not 8k  —  SID still required")
    _add_sub(s, "pair_acc = P(recon_score > paired real) on 118 pairs. A-grid has no SID.")
    _pic(s, figs["i2i_table"], 1.12, 3.55)
    # 2026-09-01, tianqi, i2i A1/A2/A3/D6 results on team slide 5
    _bullets(
        s,
        [
            "A2 i2i-only official 0.443 / pair 0.42: 59 scenes overfit. Whole-image i2i cannot replace SID.",
            "D6 keeps the SID mix and only adds 118 i2i fakes: pair_acc 0.805, official 0.977. Contest submit stays last-4 / last4+D3.",
            "SID tampered (label=2) is a local edit, not whole-image i2i. Do not mine DALL·E / EvalGEN into train.",
        ],
        0.5,
        4.85,
        12.3,
        1.9,
        15,
    )
    _notes(
        s,
        "i2i 必须是整图 triplet。A2 只训 59 组会过拟合，官方 0.44。没有 SID 的 A1/A3 官方最高 0.79。D6 在 D5 上加 118 张 fake，pair_acc 到 0.805，官方仍 0.977，不能换提交。GAN 不进训：官方和 EvalGEN 都不是 StyleGAN。",
    )
    # end
    _footer(s, 5)

    # 6 transforms
    s = prs.slides.add_slide(blank)
    _add_title(s, "Transforms: one spec for train aug and frozen eval")
    _add_sub(s, "Pillow + NumPy spec. Crop 80% is not resized back. Jitter is sampled ±20% per image.")
    _pic(s, figs["transforms"], 1.15, 4.0)
    _bullets(
        s,
        [
            "Train: online random official ops, p_clean=0.2. Eval: 14 keys frozen, seed = sha1(image|key|v2).",
            "Weakest remaining keys for fuse: JPEG-30 and resize ×0.25 (~0.984) — still far above CIFAKE.",
        ],
        0.5,
        5.3,
        12.3,
        1.6,
        16,
    )
    _notes(
        s,
        "变换不是事后补丁：训练就在同样 14 档上随机增强。CIFAKE 那条没做增强，blur/resize 会塌。评测 crop 不放大，和题面 resize then upscale 区分开。",
    )
    _footer(s, 6)

    # 7 mlp fft
    s = prs.slides.add_slide(blank)
    _add_title(s, "Model: MLP / FFT / pixel-only are dead ends on DALL·E")
    _add_sub(s, "Same frozen CLIP-B linear unless noted. Contest score is AUROC, not SID-val Acc.")
    _bullets(
        s,
        [
            "CIFAKE CLIP-B / DINOv2 / MLP probes: official formula 0.50–0.79. Do not pick the backbone here.",
            "C-Pixel (500 ADM + 500 DDPM): full official ~0.65, 1 FP / 8838 FN — never fires on DALL·E.",
            "FFT / high-pass as a second tower (SID dual RGB+HP): no shortcut, but no gain (0.966 vs SID 0.970).",
            "Takeaway: pixel/FFT can be a thin mix-in (D3), not the detector.",
        ],
        size=18,
    )
    _notes(
        s,
        "模型线从 CIFAKE 线性头和 FFT 探针开始，官方 val 接近随机。C-Pixel 在 DALL·E 上几乎全漏。频域双路没踩捷径但也没涨。所以后来只把 ADM/DDPM 当薄 mixin，不当主模型。",
    )
    _footer(s, 7)

    # 8 model bars
    s = prs.slides.add_slide(blank)
    _add_title(s, "Unfreeze / consistency / dual / 336 — then stop stacking")
    _pic(s, figs["model_bars"], 1.05, 5.7)
    _notes(
        s,
        "一致性损失是 Acc 口径最干净的涨，但比赛看 AUC：last-4 0.990 远高于 consistency 0.965。336 在小 resize/blur 的 Acc 还掉。first-4 0.974，CLIP-L last-4 0.980，都不如 CLIP-B last-4。D3 dualbranch 0.983 赢冻结 D3，仍输 fuse。",
    )
    _footer(s, 8)

    # 9 fuse
    s = prs.slides.add_slide(blank)
    _add_title(s, "Why fuse: complementary heads, not one bigger net")
    _pic(s, figs["last4_fuse"], 0.95, 2.7)
    _pic(s, figs["fuse_why"], 3.75, 2.15)
    # 2026-09-01, tianqi, fuse ablation: last-4 ranks DALL·E, D3 recovers Nova, do not retrain last-4
    _bullets(
        s,
        [
            "last-4 ranks DALL·E (0.990) but misses Nova at 0.5 (rec 0.49). D3 is the opposite: Nova rec 0.86, official only 0.978.",
            "Average logits at inference → 0.993. Training last-4 on the mix drops official to 0.976. D5/D6 fuse 0.9927 / 0.9929 do not beat D3.",
        ],
        0.5,
        6.05,
        12.3,
        1.1,
        14,
    )
    _notes(
        s,
        "两个头互补：last-4 负责官方 DALL·E 排序，D3 补 Nova 召回。不能把 last-4 再训到 mix 上（0.976）。D5/D6 更多 mixin 的 fuse 没有超过 0.993，所以提交仍是 last4+D3。0.5 召回 fuse 只有 0.56，AUC 才是分数。",
    )
    # end
    _footer(s, 9)

    # 10 official eval
    s = prs.slides.add_slide(blank)
    _add_title(s, "Eval: official demonstration val")
    _add_sub(s, "400-subset screen (seed 0, 200/200). last-4 full 13,843 = 0.989, same ranking.")
    _pic(s, figs["eval_official"], 1.15, 5.7)
    _notes(
        s,
        "官方演示集 COCO + DALL·E，训练禁入。400 是日筛，full 验证 last-4 0.989 同序。fuse 每一档都在 0.984 以上。CIFAKE 那条在 0.55 附近，用来说明域不对。",
    )
    _footer(s, 10)

    # 11 evalgen
    s = prs.slides.add_slide(blank)
    _add_title(s, "Eval: EvalGEN hold-out  (Flux / GoT / Infinity / OmniGen / Nova)")
    _pic(s, figs["evalgen"], 1.1, 5.7)
    _notes(
        s,
        "EvalGEN 五家都没进训练。GoT/OmniGen 大家接近 1.0，不要当 unseen 代理。Nova 15 档：D3 0.9865，fuse last4+D3 0.9880。last-4 召回 0.49，D3 0.86，fuse 0.56，阈值要讲清楚。",
    )
    _footer(s, 11)

    # 12 badcase fn
    s = prs.slides.add_slide(blank)
    _add_title(s, "Bad cases: FN are often non-photoreal  —  not only “perfect fakes”")
    _add_sub(s, "400 clean @0.5: fuse 1 FP / 44 FN. last-4 1 / 60. D3 4 / 20.")
    _pic(s, figs["gallery"], 1.15, 5.75)
    # 2026-09-01, tianqi, team FN notes: comics off-target; residual photoreal DALL·E
    _notes(
        s,
        "画册里 FN 大量是漫画、插画、卡通、绘画风 DALL·E，不是社交实拍目标。SID 是照片级 FLUX，画风一偏分数会过低。平台场景（实拍社交图）同一阈值下漏检预期会少。但仍有 pred 极低的实拍风 DALL·E，不能说漏检会消失。FP 几乎没有，讲平台误杀贵。",
    )
    # end
    _footer(s, 12)

    # 13 gallery product
    s = prs.slides.add_slide(blank)
    _add_title(s, "Gallery as a tool, not a screenshot dump")
    _bullets(
        s,
        [
            "scripts/badcase_gallery.py — FP/FN split, sort by worst pred, generator + condition on each card.",
            "400-subset = one HTML. Full val = paged folder (36 thumbs/page). Open index.html locally, not Jupyter.",
            "Pairwise pages: last4 → fuse fixed 16 FN, FP stayed 1. SID-aug → last4 FN 37→60 (AUC up, Acc@0.5 down).",
            "Download outputs/tables/badcase_galleries/ + badcase_compare/. Demo video: open fuse_u4_d3_400_clean.html.",
        ],
        size=18,
    )
    _notes(
        s,
        "录视频时打开本机 HTML：先 FN 插画风，再那一张 FP。对照页证明 fuse 在修漏检而不是新制造误杀。强调 0.5 不是竞赛指标，last-4 解冻后 Acc 掉、AUC 涨。",
    )
    _footer(s, 13)

    # 14 close
    s = prs.slides.add_slide(blank)
    _add_title(s, "Submit")
    _pic(s, figs["error_loop"], 1.1, 3.4)
    _bullets(
        s,
        [
            "1 ckpt: experiments/clipb16_linear_sid_unfreeze4/ckpts/best.pt   formula 0.990",
            "2 ckpts: mean logit with clipb16_linear_sid_d3_mix   formula 0.993",
            "python predict.py <dir> out.json  —  {image_path, pred=P(AIGC)}",
        ],
        0.5,
        4.7,
        12.3,
        2.2,
        18,
    )
    _notes(s, "收束：域、架构覆盖、互补头。仓库公开，predict.py 一行。不要再念 DINOv2 0.90。")
    _footer(s, 14)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


def build_demo_pptx(figs: dict[str, Path]) -> Path:
    """Six slides for the 4 min English demo. Slide 1 keeps the team's title page."""
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    n_total = 6

    s = _blank(prs)
    _accent_bar(s)
    # 2026-08-31, tianqi, keep cover layout: Team Jambuddy separate; names are not body bullets
    _add_title(s, "Robust Detection of AI-Generated Images Under Real-World Transformations", 0.75, 26)
    _add_sub(s, "TikTok TechJam 2026", 1.92, 18)
    box = s.shapes.add_textbox(Inches(0.45), Inches(2.69), Inches(12.4), Inches(0.4))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Team Jambuddy"
    _set_run(run, 20, RGB["ink"], True)
    _bullets(
        s,
        [
            "He Tianqi    tianqi008@e.ntu.edu.sg",
            "Yu Shenli    yxy021104@gmail.com",
            "Zheng Congyun    CONGYUN001@e.ntu.edu.sg",
            "Wang Weijie    wangweijiejames@163.com",
            "Tan Linyun    tanlinyun@126.com",
        ],
        0.5,
        3.15,
        12,
        3.4,
        18,
        False,
    )
    # end
    # 2026-09-01, tianqi, English 4 min speaker notes aligned with docs/demo_script.md
    _notes(
        s,
        "0:00–0:25. Hi, we are Team Jambuddy. After a photo is shared, it is JPEG-compressed, "
        "blurred, and cropped. The contest score is 0.50 times clean AUC plus 0.50 times the mean "
        "AUC under fourteen official transforms — not accuracy at 0.5. Our detector is CLIP-B, "
        "about 86 million parameters, well under the 2B cap.",
    )
    # end
    _footer(s, 1, n_total)

    s = _blank(prs)
    _add_title(s, "Mixed data")
    _add_sub(s, "last-4 trains on SID only. We generated the families below; only the architecture mix is in the fused head.")
    _pic(s, figs["arch_table"], 1.12, 3.42)
    # 2026-09-01, tianqi, demo mix slide: generated vs submit, no D-codes
    _bullets(
        s,
        [
            "1 ckpt last-4 = SID ~140k + online aug. Extra T2I (PixArt / GPT / SDXL / nano) was built, then dropped.",
            "2 ckpts: fuse last-4 with a mixed head on WildFake SD + Flux.2 + SD3.5 + ADM/DDPM (replace equal SID FLUX).",
            "Self-built from local checkpoints and ComfyUI; we will open-source the mix. No GAN / Hunyuan.",
        ],
        0.5,
        4.85,
        12.3,
        1.9,
        15,
    )
    _notes(
        s,
        "0:25–1:10. Most of the training set is still SID_Set: about 140,000 images — 70k real photos and "
        "70k FLUX fakes. The one-checkpoint detector trains on SID only. We also generated other families "
        "from local checkpoints and ComfyUI, and we will open-source that mix. If two checkpoints are "
        "allowed, we fuse last-4 with a mixed head trained on WildFake original SD, Flux.2, SD 3.5, ADM, "
        "and DDPM — replacing an equal number of SID FLUX, not stacking. We further generated PixArt, SDXL, "
        "GPT-image, and nano banana. Those extra T2I sets did not lift the contest score, so they are not "
        "in the submitted checkpoint. No GAN: contest and EvalGEN are diffusion, flow, and autoregressive.",
    )
    # end
    _footer(s, 2, n_total)

    s = _blank(prs)
    _add_title(s, "Fuse architecture")
    _pic(s, figs["fuse_arch"], 1.05, 3.55)
    _bullets(
        s,
        [
            "CLIP-B last-4 (0.990) beats CLIP-L last-4 (0.980), ResNet SID (0.779), and DINOv2 (0.79).",
            "resize 336, RGB+freq, and first-4 all lose on the contest formula. Mean-logit fuse = 0.993.",
        ],
        0.5,
        4.8,
        12.3,
        1.8,
        16,
    )
    _notes(
        s,
        "1:10–1:55. We unfreeze the last four CLIP-B vision blocks. That scores 0.990 on the official "
        "400 screen. A larger backbone does not help: CLIP-L last-4 is 0.980, ResNet SID 0.779, "
        "DINOv2 about 0.79. Resize 336, RGB plus frequency, and unfreezing the first four all lose "
        "on the contest formula. If two checkpoints are allowed, we average logits of last-4 and "
        "a mixed-data head — fuse 0.993. One command: folder in, JSON out, pred = P(AIGC). "
        "Official val never enters training.",
    )
    _footer(s, 3, n_total)

    s = _blank(prs)
    _add_title(s, "Evaluation")
    _add_sub(s, "Official val = contest score. EvalGEN is unseen generators (never trained), Nova is the hard family.")
    _pic(s, figs["eval_full"], 1.12, 4.35)
    _bullets(
        s,
        [
            "Official val = COCO real + DALL·E Advanced. SID full 13,843 formula 0.966; last-4 full 0.989.",
            "EvalGEN (Flux / GoT / Infinity / OmniGen / Nova) tests generator shift. Nova recall: mixed 0.86 vs last-4 0.49.",
        ],
        0.5,
        5.6,
        12.3,
        1.5,
        15,
    )
    _notes(
        s,
        "1:55–2:45. Left is the contest score. Official val is COCO reals plus DALL·E Advanced. "
        "Frozen SID on the full 13,843 is 0.966; last-4 on that full set is 0.989, same ranking as "
        "the 400 screen. Right is EvalGEN: Flux, GoT, Infinity, OmniGen, Nova — never trained. "
        "Nova is the hard family. Mixed data lifts Nova recall at 0.5 from 0.49 to 0.86, and "
        "Nova AUC to 0.988. Fuse keeps last-4’s DALL·E ranking and the mixed head’s Nova AUC.",
    )
    _footer(s, 4, n_total)

    s = _blank(prs)
    _add_title(s, "Bad-case gallery")
    _add_sub(s, "HTML gallery: FP / FN split, sort by worst pred, generator + path on each card. Open locally.")
    _caption(s, "FN Badcases", 0.45, 1.18)
    _caption(s, "FP Badcases", 6.85, 1.18)
    _pic_at(s, figs["fn_badcases"], 0.35, 1.48, 6.35)
    _pic_at(s, figs["fp_badcases"], 6.75, 1.48, 6.2)
    _bullets(
        s,
        [
            # 2026-09-01, tianqi, demo FN: low FPR, comics cluster, residual photoreal
            "Fuse @0.5 on 400: 1 FP / 44 FN — almost never accuses a real photo.",
            "FN cluster: comics / anime / illustration DALL·E (pred ≈ 0.001). Target is social photos, not comics.",
            "Photoreal social AIGC should miss less at the same threshold. Residual: some photoreal DALL·E still score very low.",
            # end
        ],
        0.45,
        5.42,
        12.3,
        1.4,
        14,
    )
    # 2026-09-01, tianqi, FN cluster is non-photoreal; do not claim all misses vanish in the wild
    _notes(
        s,
        "2:45–3:25. HTML gallery: false positives versus false negatives, sorted by worst prediction. "
        "At threshold 0.5, fuse on the 400 screen is one false positive and forty-four false negatives — "
        "we almost never accuse a real photo. Many misses are non-photoreal DALL·E: comics, anime, "
        "illustration, scored near 0.001. That is outside our social-photo target, so on a photoreal "
        "social feed we expect fewer misses at the same threshold. Some remaining misses are still "
        "photoreal DALL·E, so this is not a free lunch. High AUC does not mean 0.5 is the right "
        "operating point.",
    )
    # end
    _footer(s, 5, n_total)

    s = _blank(prs)
    _add_title(s, "Robustness: 15 official conditions")
    _add_sub(s, "JPEG / blur / resize / noise / jitter / center-crop.  Weakest keys: JPEG-30 and resize ×0.25.")
    _pic(s, figs["table_15cond"], 1.15, 3.35)
    _bullets(
        s,
        [
            "Contest score = 0.50×AUC_clean + 0.50× mean AUROC of the 14 transform keys.",
            "Fuse stays ≥ 0.984 on every key. Mixed data dips on hard JPEG / tiny resize; SID is flatter but lower.",
        ],
        0.5,
        4.75,
        12.3,
        1.8,
        16,
    )
    _notes(
        s,
        "3:25–4:00. Contest formula is half clean AUC and half the mean of fourteen transform keys: "
        "JPEG, blur, resize, noise, jitter, center crop. Fuse stays at or above 0.984 on every key. "
        "Weakest keys: JPEG-30 and resize ×0.25. Mixed data dips there; SID is flatter but lower. "
        "The repository is public. Install CLIP-B, load the checkpoint, run predict.py. Thank you.",
    )
    _footer(s, 6, n_total)

    OUT_DEMO.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_DEMO)
    return OUT_DEMO


def main() -> None:
    FIG.mkdir(parents=True, exist_ok=True)
    _setup_mpl()
    figs = {
        "data_path": fig_data_path(),
        "generators": fig_generators(),
        "decisions": fig_decisions(),
        "transforms": fig_transforms(),
        "model_bars": fig_model_bars(),
        "last4_fuse": fig_last4_fuse(),
        "fuse_why": fig_fuse_why(),
        "i2i_table": fig_i2i_table(),
        "fuse_arch": fig_fuse_arch(),
        "arch_table": fig_arch_table(),
        "eval_official": fig_eval_official(),
        "eval_full": fig_eval_full(),
        "evalgen": fig_evalgen(),
        "gallery": fig_gallery(),
        "error_loop": fig_error_loop(),
        "table_15cond": fig_15cond_table(),
        "fn_badcases": FIG / "fn_badcases.png",
        "fp_badcases": FIG / "fp_badcases.png",
    }
    out = build_pptx(figs)
    demo = build_demo_pptx(figs)
    print(f"wrote {out}")
    print(f"wrote {demo}")
    for k, p in figs.items():
        print(f"  {k} {p.stat().st_size} {p}")


if __name__ == "__main__":
    main()
# end
